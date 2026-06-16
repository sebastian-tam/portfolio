#!/usr/bin/env python3
"""
Create a professional PowerPoint presentation for Brand Z Pricing Intelligence Project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

# Color palette (matching portfolio theme)
NAVY = RGBColor(19, 77, 48)      # #134D30
GREEN = RGBColor(22, 163, 74)    # #16A34A
TEAL = RGBColor(52, 211, 153)    # #34D399
WHITE = RGBColor(255, 255, 255)
LIGHT_GREY = RGBColor(245, 244, 241)
TEXT_DARK = RGBColor(30, 41, 59)

def add_title_slide(prs):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Brand Z Pricing Intelligence"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "Competitive Success Report"
    p.font.size = Pt(32)
    p.font.color.rgb = TEAL

    # Project lead
    lead_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    lead_frame = lead_box.text_frame
    p = lead_frame.paragraphs[0]
    p.text = "Project Lead: Sebastian Tam (Growth Leader & AI Architect)"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE

def add_situation_slide(prs):
    """Add Situation slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "1. SITUATION: The \"Black Box\" Challenge"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    points = [
        ("Context:", "Brand Z uses demand-responsive dynamic pricing in Bath with no public API"),
        ("The Problem:", "Prices fluctuate hourly, making it impossible to benchmark Zippe's competitiveness against regulated council tariffs using static data"),
        ("Strategic Impact:", "Without real-time pricing visibility, Zippe cannot optimize pricing strategy or respond to competitive threats")
    ]

    for i, (label, text) in enumerate(points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i] if i == 0 else text_frame.paragraphs[-1]
        p.text = f"{label} {text}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        p.level = 0

def add_task_slide(prs):
    """Add Task slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "2. TASK: Reverse-Engineering the Market"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    points = [
        ("Objective:", "Build a self-improving AI system to predict live Brand Z fares for any route and vehicle class"),
        ("KPI Target:", "Achieve a Mean Absolute Percentage Error (MAPE) of <11% to ensure strategic reliability"),
        ("Success Criteria:", "Model must be accurate enough to inform real-time pricing decisions with confidence")
    ]

    for i, (label, text) in enumerate(points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i] if i == 0 else text_frame.paragraphs[-1]
        p.text = f"{label} {text}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(12)
        p.space_after = Pt(12)

def add_action_slide(prs):
    """Add Action slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "3. ACTION: The AI & Data Sprint"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Content boxes
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    points = [
        ("Data Scale:", "Automated scraping of 5,524,786 fare observations across 1,960 unique routes over 18 days"),
        ("Technical Build:", "Developed a 46-feature XGBoost Machine Learning model"),
        ("Innovation:", "Integrated 'Phase 1+2+3 Variance Features,' including cyclical hour/day patterns and Bath Council Tariff anchors as structural predictors"),
        ("Infrastructure:", "Continuous self-improving pipeline with hourly retraining capability")
    ]

    for i, (label, text) in enumerate(points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i] if i == 0 else text_frame.paragraphs[-1]
        p.text = f"{label} {text}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)
        p.space_after = Pt(10)

def add_result_slide(prs):
    """Add Result slide with accuracy chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "4. RESULT: Strategic Market Dominance"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Accuracy Chart
    x, y, cx, cy = Inches(0.5), Inches(1.1), Inches(4.5), Inches(3.5)
    chart_data = CategoryChartData()
    chart_data.categories = ['XGBoost', 'Baseline']
    chart_data.add_series('MAPE %', (7.3, 23.9))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Model Accuracy"
    chart.has_legend = False

    # Key stats boxes
    stats = [
        ("5.5M+\nObservations", Inches(5.3), Inches(1.1)),
        ("46\nFeatures", Inches(7.0), Inches(1.1)),
        ("7.3%\nMAPE", Inches(8.6), Inches(1.1))
    ]

    for stat, x_pos, y_pos in stats:
        shape = slide.shapes.add_shape(1, x_pos, y_pos, Inches(1.1), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = GREEN
        shape.line.color.rgb = GREEN

        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = 1  # Middle
        p = text_frame.paragraphs[0]
        p.text = stat
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Key findings
    findings_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2))
    findings_frame = findings_box.text_frame
    findings_frame.word_wrap = True

    findings = [
        "• Short Trips (<5mi): Brand Z maintains a 24.6% price premium over council rates",
        "• Long-Haul (20-30mi): Brand Z aggressively undercuts the council by 33.2%",
        "• Surge Peak: 14:00 Afternoon Peak is the most expensive window (+26% above baseline)"
    ]

    for i, finding in enumerate(findings):
        if i > 0:
            findings_frame.add_paragraph()
        p = findings_frame.paragraphs[i] if i == 0 else findings_frame.paragraphs[-1]
        p.text = finding
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

def add_pricing_delta_slide(prs):
    """Add Pricing Delta slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Brand Z vs. Council Pricing Strategy"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Pricing delta chart
    x, y, cx, cy = Inches(0.8), Inches(1.2), Inches(8.4), Inches(3.8)
    chart_data = CategoryChartData()
    chart_data.categories = ['5 mi\n(Premium)', '10 mi\n(Baseline)', '20 mi\n(Undercut)', '30 mi\n(Deep Undercut)']
    chart_data.add_series('Price Ratio vs Council T1', (1.246, 1.0, 0.85, 0.668))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Brand Z Pricing Positioning by Distance"
    chart.has_legend = False

    # Insight box
    insight_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(8.4), Inches(1.5))
    insight_frame = insight_box.text_frame
    insight_frame.word_wrap = True

    p = insight_frame.paragraphs[0]
    p.text = "Strategic Insight: Brand Z uses dynamic pricing to compete on short routes (premium model) while aggressively competing on long routes (volume strategy). This dual approach maximizes revenue capture across market segments."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.font.italic = True

def add_conclusion_slide(prs):
    """Add Conclusion slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Strategic Recommendations"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(8.4), Inches(3))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    recommendations = [
        "✓ Implement dynamic pricing on short routes to defend against Brand Z premium positioning",
        "✓ Develop competitive underpricing strategy for high-volume long routes",
        "✓ Deploy real-time surge monitoring at 14:00 peak to capture demand elasticity",
        "✓ Integrate Brand Z pricing signals into Zippe's live pricing algorithm"
    ]

    for i, rec in enumerate(recommendations):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i] if i == 0 else text_frame.paragraphs[-1]
        p.text = rec
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_before = Pt(10)
        p.space_after = Pt(10)

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add slides
    add_title_slide(prs)
    add_situation_slide(prs)
    add_task_slide(prs)
    add_action_slide(prs)
    add_result_slide(prs)
    add_pricing_delta_slide(prs)
    add_conclusion_slide(prs)

    # Save
    output_path = '/Users/sebastiantam/portfolio/Brand_Z_Pricing_Intelligence.pptx'
    prs.save(output_path)
    print(f"✓ PowerPoint created: {output_path}")
    return output_path

if __name__ == '__main__':
    create_presentation()
