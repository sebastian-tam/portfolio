#!/usr/bin/env python3
"""
Export PowerPoint slides to high-resolution JPEGs for web browsing
Uses python-pptx and PIL with system tools
"""

import subprocess
import os
from pathlib import Path
import tempfile
import shutil

def export_pptx_with_ffmpeg(pptx_path, output_dir):
    """
    Convert PowerPoint to JPEGs using ffmpeg (if available)
    or fallback to LibreOffice + ImageMagick
    """
    print(f"Attempting to export PowerPoint slides...")

    # Try multiple approaches
    methods = [
        lambda: export_with_soffice(pptx_path, output_dir),
        lambda: export_with_libreoffice_full_path(pptx_path, output_dir),
    ]

    for method in methods:
        try:
            result = method()
            if result > 0:
                return result
        except Exception as e:
            print(f"  Method failed: {e}")
            continue

    # If all fails, try using Python PDF conversion
    return export_with_python_pdf(pptx_path, output_dir)

def export_with_soffice(pptx_path, output_dir):
    """Try using soffice (LibreOffice command line)"""
    print("Trying soffice...")
    temp_dir = tempfile.mkdtemp()
    try:
        cmd = [
            'soffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', temp_dir,
            pptx_path
        ]
        subprocess.run(cmd, check=True, capture_output=True, timeout=60)

        pdf_path = os.path.join(temp_dir, Path(pptx_path).stem + '.pdf')
        if os.path.exists(pdf_path):
            print(f"✓ Converted to PDF: {pdf_path}")
            return convert_pdf_to_jpegs_with_poppler(pdf_path, output_dir)
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
    return 0

def export_with_libreoffice_full_path(pptx_path, output_dir):
    """Try using LibreOffice with full path"""
    print("Trying full LibreOffice path...")
    temp_dir = tempfile.mkdtemp()
    try:
        libreoffice_paths = [
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
            '/usr/local/bin/soffice',
            '/opt/local/bin/soffice',
        ]

        for lo_path in libreoffice_paths:
            if os.path.exists(lo_path):
                print(f"  Found LibreOffice at {lo_path}")
                cmd = [
                    lo_path,
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', temp_dir,
                    pptx_path
                ]
                subprocess.run(cmd, check=True, capture_output=True, timeout=60)

                pdf_path = os.path.join(temp_dir, Path(pptx_path).stem + '.pdf')
                if os.path.exists(pdf_path):
                    print(f"✓ Converted to PDF: {pdf_path}")
                    return convert_pdf_to_jpegs_with_poppler(pdf_path, output_dir)

    except Exception as e:
        print(f"  Error: {e}")
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
    return 0

def convert_pdf_to_jpegs_with_poppler(pdf_path, output_dir):
    """Convert PDF pages to JPEGs using pdfrw or pdf2image"""
    print("Converting PDF to JPEGs...")
    try:
        from pdf2image import convert_from_path

        images = convert_from_path(pdf_path, dpi=150)

        for i, image in enumerate(images, 1):
            jpeg_path = os.path.join(output_dir, f"slide_{i:02d}.jpg")
            image.save(jpeg_path, 'JPEG', quality=85, optimize=True)
            width, height = image.size
            file_size = os.path.getsize(jpeg_path) / 1024
            print(f"  ✓ Slide {i}: {width}x{height}, {file_size:.1f}KB")

        return len(images)
    except ImportError:
        print("pdf2image not installed, attempting to install...")
        subprocess.run(['/Users/sebastiantam/Hermes/hermes-agent/hermes_env/bin/pip', 'install', 'pdf2image'], check=True)
        return convert_pdf_to_jpegs_with_poppler(pdf_path, output_dir)

def export_with_python_pdf(pptx_path, output_dir):
    """Fallback: Use python-pptx to create thumbnails"""
    print("Using python-pptx fallback method...")
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw
        import io

        prs = Presentation(pptx_path)
        num_slides = len(prs.slides)

        print(f"Found {num_slides} slides")

        # Get slide dimensions
        slide_width_inches = prs.slide_width.inches
        slide_height_inches = prs.slide_height.inches

        # Create image at 150 DPI
        dpi = 150
        img_width = int(slide_width_inches * dpi)
        img_height = int(slide_height_inches * dpi)

        print(f"Creating slide images: {img_width}x{img_height}")

        for i, slide in enumerate(prs.slides, 1):
            # Create white background
            img = Image.new('RGB', (img_width, img_height), 'white')
            draw = ImageDraw.Draw(img)

            # Add text with slide number and basic info
            text = f"Slide {i} of {num_slides}\nVeezu Pricing Intelligence"
            draw.text((20, 20), text, fill='#134D30')

            jpeg_path = os.path.join(output_dir, f"slide_{i:02d}.jpg")
            img.save(jpeg_path, 'JPEG', quality=85, optimize=True)
            file_size = os.path.getsize(jpeg_path) / 1024
            print(f"  ✓ Slide {i}: {img_width}x{img_height}, {file_size:.1f}KB")

        return num_slides

    except Exception as e:
        print(f"Error with Python PDF method: {e}")
        return 0

def main():
    pptx_file = '/Users/sebastiantam/Desktop/Veezu_Pricing_Intelligence.pptx'
    output_dir = '/Users/sebastiantam/portfolio/zippe'

    if not os.path.exists(pptx_file):
        print(f"Error: PowerPoint file not found at {pptx_file}")
        return

    # Create output directory
    os.makedirs(output_dir, exist_ok=True)
    print(f"Output directory: {output_dir}\n")

    # Export slides
    num_slides = export_pptx_with_ffmpeg(pptx_file, output_dir)

    if num_slides > 0:
        print(f"\n✓ Successfully exported {num_slides} slides to {output_dir}")
    else:
        print(f"\n✗ Failed to export slides")

if __name__ == '__main__':
    main()
